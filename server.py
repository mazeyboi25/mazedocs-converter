from __future__ import annotations

import csv
import html as html_lib
import io
import json
import os
import re
import shutil
import subprocess
import threading
import time
import uuid
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import fitz
import markdown
from bs4 import BeautifulSoup
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response
from starlette.background import BackgroundTask
from fastapi.staticfiles import StaticFiles
from openpyxl import Workbook, load_workbook
from lxml import etree
from PIL import Image
from pdf2docx import Converter as PDFToDOCXConverter
from pillow_heif import register_heif_opener
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches as PptxInches
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle

register_heif_opener()

ROOT = Path(__file__).resolve().parent
MAX_UPLOAD_BYTES = 200 * 1024 * 1024

app = FastAPI(
    title="MazeDocs V2 API",
    description="Student-focused document conversion API.",
    docs_url=None,
    redoc_url=None,
    openapi_url=None,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://mazedocs.vercel.app",
        "http://127.0.0.1:5500",
        "http://localhost:5500",
        "http://127.0.0.1:8000",
        "http://localhost:8000",
    ],
    allow_origin_regex=r"https://.*\.vercel\.app",
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],
)


ROUTES: dict[str, list[dict[str, Any]]] = {
    "pdf": [
        {"target": "docx", "description": "Preserve-layout Word document: keeps each PDF page visually intact while rebuilding its text as editable positioned Word text."},
        {"target": "pptx", "description": "One visually faithful PDF page per PowerPoint slide."},
        {"target": "txt", "description": "Plain text extracted from every PDF page."},
        {"target": "png", "description": "PNG page images packaged as a ZIP file."},
        {"target": "jpg", "description": "JPEG page images packaged as a ZIP file."},
    ],
    "docx": [
        {"target": "pdf", "description": "PDF output. Uses LibreOffice when available, with a basic fallback."},
        {"target": "txt", "description": "Plain text from the Word document."},
        {"target": "html", "description": "Simple HTML reconstructed from Word paragraphs and tables."},
    ],
    "doc": [
        {"target": "docx", "description": "Upgrade legacy Word .doc to modern .docx.", "requires_libreoffice": True},
        {"target": "pdf", "description": "Convert legacy Word .doc to PDF.", "requires_libreoffice": True},
        {"target": "txt", "description": "Extract text from legacy Word .doc.", "requires_libreoffice": True},
    ],
    "pptx": [
        {"target": "docx", "description": "Word handout with slide headings, text, tables, and available images."},
        {"target": "pdf", "description": "PDF slides. Full visual fidelity uses LibreOffice; portable mode exports a text-first PDF."},
        {"target": "txt", "description": "Plain text extracted slide by slide."},
    ],
    "ppt": [
        {"target": "docx", "description": "Convert legacy PowerPoint to a Word handout.", "requires_libreoffice": True},
        {"target": "pdf", "description": "Convert legacy PowerPoint directly to PDF.", "requires_libreoffice": True},
        {"target": "txt", "description": "Extract slide text from legacy PowerPoint.", "requires_libreoffice": True},
    ],
    "xlsx": [
        {"target": "csv", "description": "CSV export. Multi-sheet workbooks are returned as a ZIP."},
        {"target": "json", "description": "Workbook data grouped by worksheet in JSON."},
        {"target": "pdf", "description": "Spreadsheet PDF. Uses LibreOffice when available, with a basic fallback."},
    ],
    "xls": [
        {"target": "xlsx", "description": "Upgrade legacy Excel .xls to .xlsx.", "requires_libreoffice": True},
        {"target": "csv", "description": "Export legacy Excel data to CSV.", "requires_libreoffice": True},
        {"target": "json", "description": "Export legacy Excel data to JSON.", "requires_libreoffice": True},
        {"target": "pdf", "description": "Convert legacy Excel workbook to PDF.", "requires_libreoffice": True},
    ],
    "csv": [
        {"target": "xlsx", "description": "Turn CSV rows into an Excel workbook."},
        {"target": "json", "description": "Convert CSV rows into JSON objects."},
    ],
    "json": [
        {"target": "csv", "description": "Flatten common JSON records into a CSV table."},
        {"target": "xlsx", "description": "Turn common JSON records into an Excel workbook."},
    ],
    "txt": [
        {"target": "docx", "description": "Put plain text into a Word document."},
        {"target": "pdf", "description": "Turn plain text into a clean printable PDF."},
    ],
    "md": [
        {"target": "html", "description": "Render Markdown as HTML."},
        {"target": "docx", "description": "Turn Markdown content into a simple Word document."},
        {"target": "pdf", "description": "Turn Markdown text into a printable PDF."},
    ],
    "html": [
        {"target": "txt", "description": "Strip HTML into plain readable text."},
        {"target": "docx", "description": "Convert headings, paragraphs, and lists into Word."},
        {"target": "pdf", "description": "Create a basic printable PDF from page text."},
    ],
    "htm": [
        {"target": "txt", "description": "Strip HTML into plain readable text."},
        {"target": "docx", "description": "Convert headings, paragraphs, and lists into Word."},
        {"target": "pdf", "description": "Create a basic printable PDF from page text."},
    ],
    "jpg": [
        {"target": "png", "description": "Convert JPEG to PNG."},
        {"target": "webp", "description": "Convert JPEG to WEBP."},
        {"target": "pdf", "description": "Place the image into a PDF page."},
    ],
    "jpeg": [
        {"target": "png", "description": "Convert JPEG to PNG."},
        {"target": "webp", "description": "Convert JPEG to WEBP."},
        {"target": "pdf", "description": "Place the image into a PDF page."},
    ],
    "png": [
        {"target": "jpg", "description": "Convert PNG to JPEG on a white background."},
        {"target": "webp", "description": "Convert PNG to WEBP."},
        {"target": "pdf", "description": "Place the image into a PDF page."},
    ],
    "webp": [
        {"target": "jpg", "description": "Convert WEBP to JPEG."},
        {"target": "png", "description": "Convert WEBP to PNG."},
        {"target": "pdf", "description": "Place the image into a PDF page."},
    ],
    "heic": [
        {"target": "jpg", "description": "Convert HEIC phone photo to JPEG."},
        {"target": "png", "description": "Convert HEIC phone photo to PNG."},
        {"target": "pdf", "description": "Place the HEIC image into a PDF page."},
    ],
}


def find_libreoffice() -> str | None:
    candidates = [
        shutil.which("libreoffice"),
        shutil.which("soffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]

    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)

    return None


def source_extension(filename: str) -> str:
    return Path(filename).suffix.lower().lstrip(".")


def safe_stem(filename: str) -> str:
    stem = Path(filename).stem.strip() or "mazedocs-file"
    bad = '<>:"/\\|?*'
    return "".join("-" if char in bad else char for char in stem)[:100]


def mime_for(path: Path) -> str:
    mapping = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv; charset=utf-8",
        ".json": "application/json; charset=utf-8",
        ".txt": "text/plain; charset=utf-8",
        ".html": "text/html; charset=utf-8",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".zip": "application/zip",
    }
    return mapping.get(path.suffix.lower(), "application/octet-stream")


def response_for(path: Path, download_name: str | None = None) -> Response:
    if not path.exists():
        raise RuntimeError("The converter did not produce an output file.")

    name = download_name or path.name
    data = path.read_bytes()

    return Response(
        content=data,
        media_type=mime_for(path),
        headers={
            "Content-Disposition": f'attachment; filename="{name}"',
            "Cache-Control": "no-store",
        },
    )


def libreoffice_convert(input_path: Path, target: str, workdir: Path) -> Path:
    executable = find_libreoffice()
    if not executable:
        raise RuntimeError(
            "This conversion route requires LibreOffice. Install LibreOffice locally or deploy the provided Dockerfile for the full MazeDocs V2 engine."
        )

    output_dir = workdir / "libreoffice-output"
    output_dir.mkdir(parents=True, exist_ok=True)

    command = [
        executable,
        "--headless",
        "--convert-to",
        target,
        "--outdir",
        str(output_dir),
        str(input_path),
    ]

    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        timeout=120,
    )

    if completed.returncode != 0:
        raise RuntimeError(
            completed.stderr.strip()
            or completed.stdout.strip()
            or "LibreOffice conversion failed."
        )

    candidates = list(output_dir.glob(f"{input_path.stem}.*"))
    matching = [path for path in candidates if path.suffix.lower() == f".{target.lower()}"]

    if matching:
        return matching[0]

    if candidates:
        return candidates[0]

    raise RuntimeError("LibreOffice finished without creating an output file.")


def pdf_to_txt(input_path: Path, output_path: Path) -> None:
    document = fitz.open(input_path)
    parts = []

    for index, page in enumerate(document, start=1):
        text = page.get_text("text").strip()
        parts.append(f"--- PAGE {index} ---\n\n{text}")

    output_path.write_text("\n\n".join(parts), encoding="utf-8")


def _clean_pdf_font_name(name: str) -> str:
    """Turn common embedded-PDF font names into Word-friendly family names."""
    value = str(name or "").strip()

    # Remove PDF subset prefixes such as ABCDEF+Poppins-Regular.
    value = re.sub(
        r"^[A-Z]{6}\+",
        "",
        value,
    )

    replacements = {
        "ArialMT": "Arial",
        "Arial-BoldMT": "Arial",
        "Arial-ItalicMT": "Arial",
        "Arial-BoldItalicMT": "Arial",
        "TimesNewRomanPSMT": "Times New Roman",
        "TimesNewRomanPS-BoldMT": "Times New Roman",
        "TimesNewRomanPS-ItalicMT": "Times New Roman",
        "TimesNewRomanPS-BoldItalicMT": "Times New Roman",
        "Calibri-Regular": "Calibri",
        "Calibri-Bold": "Calibri",
        "Calibri-Italic": "Calibri",
        "Calibri-BoldItalic": "Calibri",
        "Poppins-Regular": "Poppins",
        "Poppins-Bold": "Poppins",
        "Poppins-Italic": "Poppins",
        "Poppins-BoldItalic": "Poppins",
        "Georgia-Bold": "Georgia",
        "Georgia-Italic": "Georgia",
        "Georgia-BoldItalic": "Georgia",
    }

    if value in replacements:
        return replacements[value]

    # Strip style suffixes while leaving the family name intact.
    for suffix in (
        "-BoldItalic",
        "-BoldOblique",
        "-SemiboldItalic",
        "-SemiBoldItalic",
        "-Semibold",
        "-SemiBold",
        "-Bold",
        "-Italic",
        "-Oblique",
    ):
        if value.endswith(suffix):
            value = value[:-len(suffix)]
            break

    if value.endswith("MT") and len(value) > 2:
        value = value[:-2]

    return value or "Arial"


def _pdf_color_hex(value: int) -> str:
    """Convert PyMuPDF's packed RGB integer to Word's RRGGBB value."""
    color = int(value or 0)

    return (
        f"{(color >> 16) & 255:02X}"
        f"{(color >> 8) & 255:02X}"
        f"{color & 255:02X}"
    )


_VML_NS = "urn:schemas-microsoft-com:vml"
_OFFICE_NS = "urn:schemas-microsoft-com:office:office"
_WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_XML_NS = "http://www.w3.org/XML/1998/namespace"


def _vml_element(
    namespace: str,
    tag: str,
    attributes: dict[str, Any] | None = None,
):
    element = etree.Element(
        f"{{{namespace}}}{tag}"
    )

    prefix_to_namespace = {
        "w": _WORD_NS,
        "r": _REL_NS,
        "v": _VML_NS,
        "o": _OFFICE_NS,
    }

    for key, value in (attributes or {}).items():
        if ":" in key:
            prefix, local_name = key.split(
                ":",
                1,
            )
            attribute_namespace = prefix_to_namespace[prefix]
            element.set(
                f"{{{attribute_namespace}}}{local_name}",
                str(value),
            )
        else:
            element.set(
                key,
                str(value),
            )

    return element


def _add_fixed_page_background(
    paragraph,
    document: Document,
    image_bytes: bytes,
    page_index: int,
    page_width: float,
    page_height: float,
) -> None:
    """
    Add the PDF page's non-text visual layer behind the editable Word text.

    The background contains the original images, borders, diagrams, table
    lines, colors, watermarks, and other graphics. PDF text is removed before
    this image is rendered so edits do not expose duplicate text underneath.
    """
    relationship_id, _ = document.part.get_or_add_image(
        io.BytesIO(image_bytes)
    )

    run = OxmlElement("w:r")
    picture = OxmlElement("w:pict")

    shape = _vml_element(
        _VML_NS,
        "rect",
        {
            "id": f"MazeDocsPageBackground{page_index + 1}",
            "stroked": "f",
            "filled": "t",
            "style": (
                "position:absolute;"
                "margin-left:0pt;"
                "margin-top:0pt;"
                f"width:{page_width:.3f}pt;"
                f"height:{page_height:.3f}pt;"
                "z-index:-251654144;"
                "mso-position-horizontal-relative:page;"
                "mso-position-vertical-relative:page;"
                "mso-wrap-distance-left:0;"
                "mso-wrap-distance-right:0;"
                "mso-wrap-distance-top:0;"
                "mso-wrap-distance-bottom:0;"
            ),
        },
    )

    image_data = _vml_element(
        _VML_NS,
        "imagedata",
        {
            "r:id": relationship_id,
            "o:title": "",
        },
    )

    shape.append(image_data)
    picture.append(shape)
    run.append(picture)
    paragraph._p.append(run)


def _add_fixed_text_line(
    document: Document,
    line_bbox,
    spans: list[dict[str, Any]],
) -> None:
    """
    Recreate one PDF text line as a lightweight, absolutely positioned
    Word paragraph frame.

    Compared with VML text boxes, paragraph frames are far cheaper for Word
    to open and render, while still keeping every PDF line at its own X/Y
    position. This avoids the block reflow/overlap problem and keeps the
    page visually much closer to the source PDF.
    """
    if not spans:
        return

    text_value = "".join(
        str(span.get("text", ""))
        for span in spans
    )

    if not text_value.strip():
        return

    x0, y0, x1, y1 = (
        float(value)
        for value in line_bbox
    )

    width_points = max(
        8.0,
        (x1 - x0) * 1.12 + 4.0,
    )

    height_points = max(
        8.0,
        (y1 - y0) * 1.35 + 2.0,
    )

    paragraph = document.add_paragraph()

    p_pr = paragraph._p.get_or_add_pPr()

    # Absolute frame coordinates are expressed in twips.
    frame = OxmlElement(
        "w:framePr"
    )
    frame.set(
        qn("w:w"),
        str(
            max(
                1,
                round(width_points * 20),
            )
        ),
    )
    frame.set(
        qn("w:h"),
        str(
            max(
                1,
                round(height_points * 20),
            )
        ),
    )
    frame.set(
        qn("w:x"),
        str(
            round(x0 * 20)
        ),
    )
    frame.set(
        qn("w:y"),
        str(
            round(max(0.0, y0 - 0.8) * 20)
        ),
    )
    frame.set(
        qn("w:hAnchor"),
        "page",
    )
    frame.set(
        qn("w:vAnchor"),
        "page",
    )
    frame.set(
        qn("w:wrap"),
        "notBeside",
    )
    frame.set(
        qn("w:hRule"),
        "exact",
    )
    frame.set(
        qn("w:xAlign"),
        "left",
    )
    frame.set(
        qn("w:yAlign"),
        "top",
    )
    frame.set(
        qn("w:anchorLock"),
        "1",
    )
    p_pr.append(
        frame
    )

    spacing = OxmlElement(
        "w:spacing"
    )
    spacing.set(
        qn("w:before"),
        "0",
    )
    spacing.set(
        qn("w:after"),
        "0",
    )
    spacing.set(
        qn("w:line"),
        str(
            max(
                1,
                round(height_points * 20)
            )
        ),
    )
    spacing.set(
        qn("w:lineRule"),
        "exact",
    )
    p_pr.append(
        spacing
    )

    ind = OxmlElement(
        "w:ind"
    )
    ind.set(
        qn("w:left"),
        "0",
    )
    ind.set(
        qn("w:right"),
        "0",
    )
    ind.set(
        qn("w:firstLine"),
        "0",
    )
    p_pr.append(
        ind
    )

    # Prevent Word from trying to keep these absolute lines together and
    # accidentally moving them to another page.
    contextual = OxmlElement(
        "w:contextualSpacing"
    )
    p_pr.append(
        contextual
    )

    for span in spans:
        span_text = str(
            span.get("text", "")
        )

        if not span_text:
            continue

        run = paragraph.add_run(
            span_text
        )

        font = run.font

        raw_font_name = str(
            span.get(
                "font",
                "Arial",
            )
        )
        font.name = _clean_pdf_font_name(
            raw_font_name
        )

        font.size = Pt(
            max(
                5.0,
                min(
                    72.0,
                    float(
                        span.get(
                            "size",
                            10.0,
                        )
                    ),
                ),
            )
        )

        font_name_lower = raw_font_name.lower()
        flags = int(
            span.get(
                "flags",
                0,
            )
        )

        font.bold = bool(
            "bold" in font_name_lower
            or
            flags & 16
        )

        font.italic = bool(
            "italic" in font_name_lower
            or
            "oblique" in font_name_lower
            or
            flags & 2
        )

        color_hex = _pdf_color_hex(
            span.get(
                "color",
                0,
            )
        )

        try:
            font.color.rgb = RGBColor.from_string(
                color_hex
            )
        except Exception:
            pass

        # Preserve spaces exactly.
        run_xml = run._r
        for text_node in run_xml.findall(
            qn("w:t")
        ):
            text_node.set(
                f"{{{_XML_NS}}}space",
                "preserve",
            )



def _page_editable_lines(page) -> list[tuple[Any, list[dict[str, Any]]]]:
    """Return visible text lines and their styled PDF spans in page order."""
    lines: list[tuple[Any, list[dict[str, Any]]]] = []

    page_data = page.get_text(
        "dict",
        sort=True,
    )

    for block in page_data.get(
        "blocks",
        [],
    ):
        if int(block.get("type", 0)) != 0:
            continue

        for line in block.get(
            "lines",
            [],
        ):
            spans = list(
                line.get(
                    "spans",
                    [],
                )
            )

            if not spans:
                continue

            line_text = "".join(
                str(span.get("text", ""))
                for span in spans
            )

            if not line_text.strip():
                continue

            lines.append((
                line.get(
                    "bbox",
                    (0, 0, 0, 0),
                ),
                spans,
            ))

    return lines


def _render_non_text_page_layer(
    source_document,
    page_index: int,
    lines: list[tuple[Any, list[dict[str, Any]]]],
) -> bytes:
    """
    Render one PDF page after removing its text objects.

    Redactions use fill=None and leave images / vector graphics untouched, so
    the page's design survives while the editable Word text can sit on top.
    """
    visual_document = fitz.open()
    visual_document.insert_pdf(
        source_document,
        from_page=page_index,
        to_page=page_index,
    )

    visual_page = visual_document[0]

    for bbox, _ in lines:
        rectangle = fitz.Rect(
            bbox
        )

        # A tiny expansion catches anti-aliased glyph edges without erasing
        # nearby table borders or design elements.
        rectangle.x0 -= 0.25
        rectangle.y0 -= 0.20
        rectangle.x1 += 0.25
        rectangle.y1 += 0.20

        visual_page.add_redact_annot(
            rectangle,
            fill=None,
            cross_out=False,
        )

    if lines:
        visual_page.apply_redactions(
            images=0,
            graphics=0,
            text=0,
        )

    pixmap = visual_page.get_pixmap(
        matrix=fitz.Matrix(
            1.30,
            1.30,
        ),
        alpha=False,
    )

    image_bytes = pixmap.tobytes(
        "jpeg",
        jpg_quality=84,
    )

    visual_document.close()

    return image_bytes


def _preserve_layout_pdf_to_docx(
    input_path: Path,
    output_path: Path,
) -> None:
    """
    Preserve-layout PDF -> DOCX conversion.

    Each PDF page becomes exactly one Word page.
    - Non-text visuals stay as the same high-quality fixed background layer.
    - Every PDF text line is placed at its original page coordinates.
    - Text remains editable.
    - Native Word paragraph frames replace the old VML text boxes, which
      greatly reduces document-open overhead without grouping lines together.
    """
    source = fitz.open(
        input_path
    )

    try:
        if len(source) == 0:
            raise RuntimeError(
                "The PDF contains no pages."
            )

        document = Document()

        # Remove the default body paragraph so it cannot introduce visible
        # layout drift on the first page.
        if document.paragraphs:
            default_paragraph = document.paragraphs[0]
            default_paragraph._element.getparent().remove(
                default_paragraph._element
            )

        first_page = source[0]
        first_section = document.sections[0]
        first_section.page_width = Pt(
            first_page.rect.width
        )
        first_section.page_height = Pt(
            first_page.rect.height
        )
        first_section.left_margin = Pt(0)
        first_section.right_margin = Pt(0)
        first_section.top_margin = Pt(0)
        first_section.bottom_margin = Pt(0)
        first_section.header_distance = Pt(0)
        first_section.footer_distance = Pt(0)

        for page_index, page in enumerate(source):
            section = document.sections[-1]
            section.page_width = Pt(
                page.rect.width
            )
            section.page_height = Pt(
                page.rect.height
            )
            section.left_margin = Pt(0)
            section.right_margin = Pt(0)
            section.top_margin = Pt(0)
            section.bottom_margin = Pt(0)
            section.header_distance = Pt(0)
            section.footer_distance = Pt(0)

            editable_lines = _page_editable_lines(
                page
            )

            background_bytes = _render_non_text_page_layer(
                source,
                page_index,
                editable_lines,
            )

            # One lightweight paragraph owns the full-page visual background.
            background_paragraph = document.add_paragraph()
            background_paragraph.paragraph_format.space_before = Pt(0)
            background_paragraph.paragraph_format.space_after = Pt(0)
            background_paragraph.paragraph_format.line_spacing = Pt(1)

            _add_fixed_page_background(
                background_paragraph,
                document,
                background_bytes,
                page_index,
                page.rect.width,
                page.rect.height,
            )

            # Keep every PDF line independent so Word cannot reflow a whole
            # block and destroy columns, bullets, sidebars, or table content.
            for bbox, spans in editable_lines:
                _add_fixed_text_line(
                    document,
                    bbox,
                    spans,
                )

            if page_index < len(source) - 1:
                page_break_paragraph = document.add_paragraph()
                page_break_paragraph.paragraph_format.space_before = Pt(0)
                page_break_paragraph.paragraph_format.space_after = Pt(0)

                run = page_break_paragraph.add_run()
                run.add_break()

                break_node = run._r[-1]
                break_node.set(
                    qn("w:type"),
                    "page",
                )

        document.save(
            output_path
        )

    finally:
        source.close()



def pdf_to_docx(
    input_path: Path,
    output_path: Path,
) -> None:
    """
    MazeDocs PDF -> Word preserve-layout route.

    This route favors one-to-one page fidelity and editable positioned text
    instead of normal Word reflow. It is intended for reviewers, handouts,
    forms, certificates, designed notes, and other PDFs whose original layout
    matters.
    """
    try:
        _preserve_layout_pdf_to_docx(
            input_path,
            output_path,
        )
    except RuntimeError:
        raise
    except Exception as error:
        raise RuntimeError(
            f"PDF to Word conversion failed: {error}"
        ) from error

    if (
        not output_path.exists()
        or
        output_path.stat().st_size == 0
    ):
        raise RuntimeError(
            "PDF to Word conversion did not create a valid DOCX file."
        )


def pdf_to_pptx(input_path: Path, output_path: Path, workdir: Path) -> None:
    source = fitz.open(input_path)
    presentation = Presentation()

    first_rect = source[0].rect if len(source) else fitz.Rect(0, 0, 595, 842)
    ratio = first_rect.width / first_rect.height if first_rect.height else 0.707

    presentation.slide_width = PptxInches(10)
    presentation.slide_height = PptxInches(10 / ratio)

    blank_layout = presentation.slide_layouts[6]

    for index, page in enumerate(source, start=1):
        pixmap = page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
        image_path = workdir / f"page-{index:03d}.png"
        pixmap.save(image_path)

        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    presentation.save(output_path)


def pdf_to_image_zip(input_path: Path, output_path: Path, image_format: str, workdir: Path) -> None:
    source = fitz.open(input_path)
    image_paths: list[Path] = []

    for index, page in enumerate(source, start=1):
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        path = workdir / f"page-{index:03d}.{image_format}"

        if image_format == "png":
            pixmap.save(path)
        else:
            png_bytes = pixmap.tobytes("png")
            with Image.open(io.BytesIO(png_bytes)) as image:
                image.convert("RGB").save(path, format="JPEG", quality=90, optimize=True)

        image_paths.append(path)

    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for path in image_paths:
            archive.write(path, arcname=path.name)


def docx_to_txt(input_path: Path, output_path: Path) -> None:
    document = Document(input_path)
    parts = [paragraph.text for paragraph in document.paragraphs]

    for table in document.tables:
        parts.append("")
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))

    output_path.write_text("\n".join(parts), encoding="utf-8")


def docx_to_html(input_path: Path, output_path: Path) -> None:
    document = Document(input_path)
    pieces = ["<!doctype html><html><head><meta charset='utf-8'><title>MazeDocs Export</title></head><body>"]

    for paragraph in document.paragraphs:
        text = html_lib.escape(paragraph.text)
        style_name = (paragraph.style.name or "").lower()

        if not text.strip():
            continue

        if style_name.startswith("heading"):
            level = 2
            try:
                level = max(1, min(6, int(style_name.split()[-1])))
            except Exception:
                pass
            pieces.append(f"<h{level}>{text}</h{level}>")
        else:
            pieces.append(f"<p>{text}</p>")

    for table in document.tables:
        pieces.append("<table border='1' cellspacing='0' cellpadding='6'>")
        for row in table.rows:
            pieces.append("<tr>" + "".join(f"<td>{html_lib.escape(cell.text)}</td>" for cell in row.cells) + "</tr>")
        pieces.append("</table>")

    pieces.append("</body></html>")
    output_path.write_text("\n".join(pieces), encoding="utf-8")


def pptx_text_lines(input_path: Path) -> list[str]:
    presentation = Presentation(input_path)
    output: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        output.append(f"--- SLIDE {slide_number} ---")

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    output.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    output.append("\t".join(cell.text for cell in row.cells))

        output.append("")

    return output


def pptx_to_txt(input_path: Path, output_path: Path) -> None:
    output_path.write_text("\n".join(pptx_text_lines(input_path)), encoding="utf-8")


def pptx_to_docx(input_path: Path, output_path: Path) -> None:
    presentation = Presentation(input_path)
    document = Document()
    document.add_heading(input_path.stem, level=0)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        document.add_heading(f"Slide {slide_number}", level=1)

        seen_text: set[str] = set()

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text and text not in seen_text:
                    document.add_paragraph(text)
                    seen_text.add(text)

            if getattr(shape, "has_table", False):
                source_table = shape.table
                word_table = document.add_table(
                    rows=len(source_table.rows),
                    cols=len(source_table.columns),
                )
                word_table.style = "Table Grid"

                for row_index, row in enumerate(source_table.rows):
                    for col_index, cell in enumerate(row.cells):
                        word_table.cell(row_index, col_index).text = cell.text

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    document.add_picture(io.BytesIO(shape.image.blob), width=Inches(5.8))
            except Exception:
                pass

        if slide_number < len(presentation.slides):
            document.add_page_break()

    document.save(output_path)


def text_to_docx(text: str, output_path: Path, title: str = "MazeDocs Export") -> None:
    document = Document()
    document.add_heading(title, level=0)

    for line in text.splitlines():
        if line.strip():
            document.add_paragraph(line)
        else:
            document.add_paragraph("")

    document.save(output_path)


def text_to_pdf(text: str, output_path: Path, title: str = "MazeDocs Export") -> None:
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "MazeDocsBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        spaceAfter=6,
    )

    document = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=18 * mm,
        leftMargin=18 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
        title=title,
    )

    story = [Paragraph(html_lib.escape(title), styles["Title"]), Spacer(1, 8)]

    for line in text.splitlines():
        escaped = html_lib.escape(line) or "&nbsp;"
        story.append(Paragraph(escaped, body))

    document.build(story)


def basic_docx_to_pdf(input_path: Path, output_path: Path) -> None:
    document = Document(input_path)
    lines = [paragraph.text for paragraph in document.paragraphs]

    for table in document.tables:
        lines.append("")
        lines.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)

    text_to_pdf("\n".join(lines), output_path, input_path.stem)


def basic_pptx_to_pdf(input_path: Path, output_path: Path) -> None:
    styles = getSampleStyleSheet()
    document = SimpleDocTemplate(str(output_path), pagesize=landscape(A4), title=input_path.stem)
    story = []

    presentation = Presentation(input_path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        story.append(Paragraph(f"Slide {slide_number}", styles["Heading1"]))

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                story.append(Paragraph(html_lib.escape(shape.text).replace("\n", "<br/>"), styles["BodyText"]))
                story.append(Spacer(1, 5))

        if slide_number < len(presentation.slides):
            story.append(PageBreak())

    document.build(story)


def xlsx_to_json(input_path: Path, output_path: Path) -> None:
    workbook = load_workbook(input_path, data_only=True, read_only=True)
    result: dict[str, list[dict[str, Any]]] = {}

    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            result[sheet.title] = []
            continue

        headers = [str(value) if value is not None else f"column_{index + 1}" for index, value in enumerate(rows[0])]
        records = []

        for row in rows[1:]:
            records.append({headers[index]: row[index] if index < len(row) else None for index in range(len(headers))})

        result[sheet.title] = records

    output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False, default=str), encoding="utf-8")


def xlsx_to_csv(input_path: Path, output_path: Path, workdir: Path) -> Path:
    workbook = load_workbook(input_path, data_only=True, read_only=True)

    if len(workbook.worksheets) == 1:
        sheet = workbook.worksheets[0]
        with output_path.open("w", encoding="utf-8-sig", newline="") as handle:
            writer = csv.writer(handle)
            for row in sheet.iter_rows(values_only=True):
                writer.writerow(["" if value is None else value for value in row])
        return output_path

    zip_path = output_path.with_suffix(".zip")

    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for sheet in workbook.worksheets:
            csv_path = workdir / f"{sheet.title[:50] or 'sheet'}.csv"
            with csv_path.open("w", encoding="utf-8-sig", newline="") as handle:
                writer = csv.writer(handle)
                for row in sheet.iter_rows(values_only=True):
                    writer.writerow(["" if value is None else value for value in row])
            archive.write(csv_path, arcname=csv_path.name)

    return zip_path


def xlsx_to_basic_pdf(input_path: Path, output_path: Path) -> None:
    workbook = load_workbook(input_path, data_only=True, read_only=True)
    styles = getSampleStyleSheet()
    story = []

    for sheet_index, sheet in enumerate(workbook.worksheets):
        story.append(Paragraph(html_lib.escape(sheet.title), styles["Heading1"]))
        rows = [["" if value is None else str(value) for value in row] for row in sheet.iter_rows(values_only=True)]

        if rows:
            max_columns = min(max(len(row) for row in rows), 10)
            trimmed = [row[:max_columns] for row in rows[:150]]
            table = Table(trimmed, repeatRows=1)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#d7dcff")),
                ("FONTSIZE", (0, 0), (-1, -1), 6),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(table)

        if sheet_index < len(workbook.worksheets) - 1:
            story.append(PageBreak())

    document = SimpleDocTemplate(str(output_path), pagesize=landscape(A4), title=input_path.stem)
    document.build(story)


def csv_to_xlsx(input_path: Path, output_path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"

    with input_path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            sheet.append(row)

    workbook.save(output_path)


def csv_to_json(input_path: Path, output_path: Path) -> None:
    with input_path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        records = list(reader)

    output_path.write_text(json.dumps(records, indent=2, ensure_ascii=False), encoding="utf-8")


def normalize_json_records(data: Any) -> list[dict[str, Any]]:
    if isinstance(data, list):
        if all(isinstance(item, dict) for item in data):
            return data
        return [{"value": item} for item in data]

    if isinstance(data, dict):
        for value in data.values():
            if isinstance(value, list) and all(isinstance(item, dict) for item in value):
                return value
        return [data]

    return [{"value": data}]


def json_to_csv(input_path: Path, output_path: Path) -> None:
    data = json.loads(input_path.read_text(encoding="utf-8-sig"))
    records = normalize_json_records(data)
    headers: list[str] = []

    for record in records:
        for key in record.keys():
            if key not in headers:
                headers.append(key)

    with output_path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        for record in records:
            writer.writerow({key: json.dumps(record.get(key), ensure_ascii=False) if isinstance(record.get(key), (dict, list)) else record.get(key) for key in headers})


def json_to_xlsx(input_path: Path, output_path: Path) -> None:
    data = json.loads(input_path.read_text(encoding="utf-8-sig"))
    records = normalize_json_records(data)
    headers: list[str] = []

    for record in records:
        for key in record.keys():
            if key not in headers:
                headers.append(key)

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(headers)

    for record in records:
        sheet.append([
            json.dumps(record.get(key), ensure_ascii=False) if isinstance(record.get(key), (dict, list)) else record.get(key)
            for key in headers
        ])

    workbook.save(output_path)


def image_convert(input_path: Path, output_path: Path, target: str) -> None:
    with Image.open(input_path) as image:
        if target == "pdf":
            if image.mode in ("RGBA", "LA"):
                canvas = Image.new("RGB", image.size, "white")
                alpha = image.getchannel("A") if "A" in image.getbands() else None
                canvas.paste(image.convert("RGB"), mask=alpha)
                image = canvas
            else:
                image = image.convert("RGB")
            image.save(output_path, "PDF", resolution=150)
            return

        if target in {"jpg", "jpeg"}:
            if image.mode in ("RGBA", "LA"):
                background = Image.new("RGB", image.size, "white")
                alpha = image.getchannel("A") if "A" in image.getbands() else None
                background.paste(image.convert("RGB"), mask=alpha)
                image = background
            else:
                image = image.convert("RGB")
            image.save(output_path, "JPEG", quality=92, optimize=True)
            return

        if target == "png":
            image.save(output_path, "PNG", optimize=True)
            return

        if target == "webp":
            image.save(output_path, "WEBP", quality=90, method=6)
            return

    raise RuntimeError("Unsupported image target.")


def html_to_text_and_docx(input_path: Path, target: str, output_path: Path) -> None:
    raw = input_path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")

    if target == "txt":
        output_path.write_text(soup.get_text("\n", strip=True), encoding="utf-8")
        return

    document = Document()

    for node in soup.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
        text = node.get_text(" ", strip=True)
        if not text:
            continue
        if node.name and node.name.startswith("h"):
            document.add_heading(text, level=min(int(node.name[1]), 4))
        else:
            document.add_paragraph(text, style="List Bullet" if node.name == "li" else None)

    document.save(output_path)


def convert_file(input_path: Path, source: str, target: str, workdir: Path) -> Path:
    base = safe_stem(input_path.name)
    output_path = workdir / f"{base}-mazedocs.{target}"

    if source == "pdf":
        if target == "docx":
            pdf_to_docx(input_path, output_path)
        elif target == "pptx":
            pdf_to_pptx(input_path, output_path, workdir)
        elif target == "txt":
            pdf_to_txt(input_path, output_path)
        elif target in {"png", "jpg"}:
            output_path = workdir / f"{base}-pages.zip"
            pdf_to_image_zip(input_path, output_path, target, workdir)
        else:
            raise RuntimeError("Unsupported PDF conversion.")
        return output_path

    if source == "docx":
        if target == "txt":
            docx_to_txt(input_path, output_path)
        elif target == "html":
            docx_to_html(input_path, output_path)
        elif target == "pdf":
            if find_libreoffice():
                output_path = libreoffice_convert(input_path, "pdf", workdir)
            else:
                basic_docx_to_pdf(input_path, output_path)
        else:
            raise RuntimeError("Unsupported DOCX conversion.")
        return output_path

    if source == "doc":
        if not find_libreoffice():
            raise RuntimeError("Legacy .doc conversion requires LibreOffice.")
        if target == "docx":
            return libreoffice_convert(input_path, "docx", workdir)
        if target == "pdf":
            return libreoffice_convert(input_path, "pdf", workdir)
        if target == "txt":
            intermediate = libreoffice_convert(input_path, "docx", workdir)
            docx_to_txt(intermediate, output_path)
            return output_path

    if source == "pptx":
        if target == "docx":
            pptx_to_docx(input_path, output_path)
        elif target == "txt":
            pptx_to_txt(input_path, output_path)
        elif target == "pdf":
            if find_libreoffice():
                output_path = libreoffice_convert(input_path, "pdf", workdir)
            else:
                basic_pptx_to_pdf(input_path, output_path)
        else:
            raise RuntimeError("Unsupported PPTX conversion.")
        return output_path

    if source == "ppt":
        if not find_libreoffice():
            raise RuntimeError("Legacy .ppt conversion requires LibreOffice.")
        if target == "pdf":
            return libreoffice_convert(input_path, "pdf", workdir)
        intermediate = libreoffice_convert(input_path, "pptx", workdir)
        if target == "docx":
            pptx_to_docx(intermediate, output_path)
            return output_path
        if target == "txt":
            pptx_to_txt(intermediate, output_path)
            return output_path

    if source == "xlsx":
        if target == "json":
            xlsx_to_json(input_path, output_path)
        elif target == "csv":
            output_path = xlsx_to_csv(input_path, output_path, workdir)
        elif target == "pdf":
            if find_libreoffice():
                output_path = libreoffice_convert(input_path, "pdf", workdir)
            else:
                xlsx_to_basic_pdf(input_path, output_path)
        else:
            raise RuntimeError("Unsupported XLSX conversion.")
        return output_path

    if source == "xls":
        if not find_libreoffice():
            raise RuntimeError("Legacy .xls conversion requires LibreOffice.")
        if target == "xlsx":
            return libreoffice_convert(input_path, "xlsx", workdir)
        if target == "pdf":
            return libreoffice_convert(input_path, "pdf", workdir)
        intermediate = libreoffice_convert(input_path, "xlsx", workdir)
        if target == "json":
            xlsx_to_json(intermediate, output_path)
            return output_path
        if target == "csv":
            return xlsx_to_csv(intermediate, output_path, workdir)

    if source == "csv":
        if target == "xlsx":
            csv_to_xlsx(input_path, output_path)
        elif target == "json":
            csv_to_json(input_path, output_path)
        else:
            raise RuntimeError("Unsupported CSV conversion.")
        return output_path

    if source == "json":
        if target == "csv":
            json_to_csv(input_path, output_path)
        elif target == "xlsx":
            json_to_xlsx(input_path, output_path)
        else:
            raise RuntimeError("Unsupported JSON conversion.")
        return output_path

    if source == "txt":
        text = input_path.read_text(encoding="utf-8", errors="replace")
        if target == "docx":
            text_to_docx(text, output_path, input_path.stem)
        elif target == "pdf":
            text_to_pdf(text, output_path, input_path.stem)
        else:
            raise RuntimeError("Unsupported text conversion.")
        return output_path

    if source == "md":
        text = input_path.read_text(encoding="utf-8", errors="replace")
        if target == "html":
            body = markdown.markdown(text, extensions=["tables", "fenced_code"])
            output_path.write_text(f"<!doctype html><html><head><meta charset='utf-8'></head><body>{body}</body></html>", encoding="utf-8")
        elif target == "docx":
            temp_html = workdir / "markdown.html"
            body = markdown.markdown(text, extensions=["tables", "fenced_code"])
            temp_html.write_text(f"<html><body>{body}</body></html>", encoding="utf-8")
            html_to_text_and_docx(temp_html, "docx", output_path)
        elif target == "pdf":
            text_to_pdf(text, output_path, input_path.stem)
        else:
            raise RuntimeError("Unsupported Markdown conversion.")
        return output_path

    if source in {"html", "htm"}:
        if target in {"txt", "docx"}:
            html_to_text_and_docx(input_path, target, output_path)
        elif target == "pdf":
            soup = BeautifulSoup(input_path.read_text(encoding="utf-8", errors="replace"), "html.parser")
            text_to_pdf(soup.get_text("\n", strip=True), output_path, input_path.stem)
        else:
            raise RuntimeError("Unsupported HTML conversion.")
        return output_path

    if source in {"jpg", "jpeg", "png", "webp", "heic"}:
        image_convert(input_path, output_path, target)
        return output_path

    raise RuntimeError("This conversion route is not implemented.")


def health_payload() -> dict[str, Any]:
    executable = find_libreoffice()

    return {
        "ok": True,
        "service": "MazeDocs V2",
        "engine": "python-converter",
        "libreoffice_available": bool(executable),
        "libreoffice_path": executable,
        "max_upload_bytes": MAX_UPLOAD_BYTES,
        "conversion_mode": "background-jobs",
        "job_api": True,
        "supported_inputs": sorted(ROUTES.keys()),
    }


# ============================================================
# BACKGROUND CONVERSION JOBS
# ============================================================

JOB_RETENTION_SECONDS = 2 * 60 * 60
JOBS: dict[str, dict[str, Any]] = {}
JOBS_LOCK = threading.Lock()


def _remove_job(job_id: str) -> None:
    with JOBS_LOCK:
        job = JOBS.pop(job_id, None)

    if not job:
        return

    workdir = Path(job.get("workdir", ""))

    if workdir.exists():
        shutil.rmtree(
            workdir,
            ignore_errors=True,
        )


def _cleanup_old_jobs() -> None:
    now = time.time()
    stale: list[str] = []

    with JOBS_LOCK:
        for job_id, job in JOBS.items():
            if job.get("status") == "processing":
                continue

            created_at = float(
                job.get("created_at", now)
            )

            if (
                now - created_at
                >
                JOB_RETENTION_SECONDS
            ):
                stale.append(job_id)

    for job_id in stale:
        _remove_job(job_id)


def _run_conversion_job(job_id: str) -> None:
    """
    Run a conversion after the upload request has already returned.

    This is intentionally separate from the upload request so Railway does
    not see one long silent HTTP request while pdf2docx / LibreOffice works.
    """
    with JOBS_LOCK:
        job = JOBS.get(job_id)

        if not job:
            return

        job["status"] = "processing"
        job["stage"] = "converting"
        job["started_at"] = time.time()

        input_path = Path(job["input_path"])
        workdir = Path(job["workdir"])
        source = str(job["source"])
        target = str(job["target"])
        original_name = str(job["filename"])

    try:
        output_path = convert_file(
            input_path,
            source,
            target,
            workdir,
        )

        if (
            not output_path.exists()
            or
            output_path.stat().st_size == 0
        ):
            raise RuntimeError(
                "The converter did not produce a valid output file."
            )

        final_name = output_path.name

        if output_path.stem.startswith("source"):
            final_name = output_path.name.replace(
                "source",
                safe_stem(original_name),
                1,
            )

        with JOBS_LOCK:
            current = JOBS.get(job_id)

            if current is not None:
                current.update({
                    "status": "done",
                    "stage": "ready",
                    "finished_at": time.time(),
                    "output_path": str(output_path),
                    "download_name": final_name,
                    "output_size": output_path.stat().st_size,
                })

    except Exception as error:
        with JOBS_LOCK:
            current = JOBS.get(job_id)

            if current is not None:
                current.update({
                    "status": "error",
                    "stage": "failed",
                    "finished_at": time.time(),
                    "error": str(error),
                })


async def _save_job_upload(
    file: UploadFile,
    target: str,
) -> str:
    filename = file.filename or "upload"
    source = source_extension(filename)
    target = target.lower().strip().lstrip(".")

    if source not in ROUTES:
        raise HTTPException(
            400,
            f".{source or '?'} is not supported.",
        )

    valid_targets = {
        route["target"]
        for route in ROUTES[source]
    }

    if target not in valid_targets:
        raise HTTPException(
            400,
            f"Cannot convert .{source} to .{target}.",
        )

    route = next(
        route
        for route in ROUTES[source]
        if route["target"] == target
    )

    if (
        route.get("requires_libreoffice")
        and
        not find_libreoffice()
    ):
        raise HTTPException(
            422,
            "This conversion route requires LibreOffice.",
        )

    _cleanup_old_jobs()

    job_id = uuid.uuid4().hex
    workdir = Path(
        tempfile.mkdtemp(
            prefix=f"mazedocs-job-{job_id[:8]}-"
        )
    )
    input_path = workdir / f"source.{source}"

    received_bytes = 0
    chunk_size = 1024 * 1024

    try:
        with input_path.open("wb") as destination:
            while True:
                chunk = await file.read(chunk_size)

                if not chunk:
                    break

                received_bytes += len(chunk)

                if received_bytes > MAX_UPLOAD_BYTES:
                    raise HTTPException(
                        413,
                        (
                            "This file is larger than the MazeDocs "
                            f"{MAX_UPLOAD_BYTES // (1024 * 1024)} MB "
                            "application limit."
                        ),
                    )

                destination.write(chunk)

        if received_bytes == 0:
            raise HTTPException(
                400,
                "The uploaded file is empty.",
            )

        with JOBS_LOCK:
            JOBS[job_id] = {
                "id": job_id,
                "status": "queued",
                "stage": "queued",
                "created_at": time.time(),
                "filename": filename,
                "source": source,
                "target": target,
                "size": received_bytes,
                "workdir": str(workdir),
                "input_path": str(input_path),
            }

        return job_id

    except Exception:
        shutil.rmtree(
            workdir,
            ignore_errors=True,
        )
        raise

    finally:
        await file.close()


async def convert_upload(
    file: UploadFile,
    target: str,
) -> FileResponse:
    """
    Stream an uploaded file to temporary disk, convert it, then stream the
    generated file back to the browser.

    This avoids loading a 50 MB / 100 MB / 200 MB source file and its output
    into Python memory at the same time.
    """
    filename = file.filename or "upload"
    source = source_extension(filename)
    target = target.lower().strip().lstrip(".")

    if source not in ROUTES:
        raise HTTPException(
            400,
            f".{source or '?'} is not supported.",
        )

    valid_targets = {
        route["target"]
        for route in ROUTES[source]
    }

    if target not in valid_targets:
        raise HTTPException(
            400,
            f"Cannot convert .{source} to .{target}.",
        )

    route = next(
        route
        for route in ROUTES[source]
        if route["target"] == target
    )

    if (
        route.get("requires_libreoffice")
        and
        not find_libreoffice()
    ):
        raise HTTPException(
            422,
            "This conversion route requires LibreOffice.",
        )

    workdir = Path(
        tempfile.mkdtemp(
            prefix="mazedocs-"
        )
    )

    input_path = (
        workdir
        /
        f"source.{source}"
    )

    received_bytes = 0
    chunk_size = 1024 * 1024

    try:
        with input_path.open("wb") as destination:
            while True:
                chunk = await file.read(
                    chunk_size
                )

                if not chunk:
                    break

                received_bytes += len(
                    chunk
                )

                if (
                    received_bytes
                    >
                    MAX_UPLOAD_BYTES
                ):
                    raise HTTPException(
                        413,
                        (
                            "This file is larger than the MazeDocs "
                            f"{MAX_UPLOAD_BYTES // (1024 * 1024)} MB "
                            "application limit."
                        ),
                    )

                destination.write(
                    chunk
                )

        if received_bytes == 0:
            raise HTTPException(
                400,
                "The uploaded file is empty.",
            )

        try:
            output_path = convert_file(
                input_path,
                source,
                target,
                workdir,
            )
        except RuntimeError as error:
            raise HTTPException(
                422,
                str(error),
            ) from error
        except HTTPException:
            raise
        except Exception as error:
            raise HTTPException(
                500,
                f"Conversion failed: {error}",
            ) from error

        if (
            not output_path.exists()
            or
            output_path.stat().st_size == 0
        ):
            raise HTTPException(
                500,
                "The converter did not produce a valid output file.",
            )

        final_name = output_path.name

        if output_path.stem.startswith(
            "source"
        ):
            final_name = (
                output_path.name.replace(
                    "source",
                    safe_stem(
                        filename
                    ),
                    1,
                )
            )

        cleanup = BackgroundTask(
            shutil.rmtree,
            workdir,
            ignore_errors=True,
        )

        return FileResponse(
            path=output_path,
            media_type=mime_for(
                output_path
            ),
            filename=final_name,
            headers={
                "Cache-Control":
                    "no-store",
            },
            background=cleanup,
        )

    except Exception:
        shutil.rmtree(
            workdir,
            ignore_errors=True,
        )
        raise

    finally:
        await file.close()


@app.post("/api/jobs", status_code=202)
async def create_conversion_job(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    target: str = Form(...),
) -> dict[str, Any]:
    """
    Upload the source file, return immediately with a job id, and run the
    expensive conversion after the HTTP upload request has finished.
    """
    job_id = await _save_job_upload(
        file,
        target,
    )

    background_tasks.add_task(
        _run_conversion_job,
        job_id,
    )

    return {
        "ok": True,
        "job_id": job_id,
        "status": "queued",
        "status_url": f"/api/jobs/{job_id}",
        "download_url": f"/api/jobs/{job_id}/download",
    }


@app.get("/api/jobs/{job_id}")
def conversion_job_status(
    job_id: str,
) -> dict[str, Any]:
    _cleanup_old_jobs()

    with JOBS_LOCK:
        job = JOBS.get(job_id)

        if not job:
            raise HTTPException(
                404,
                "Conversion job was not found or has expired.",
            )

        payload = {
            "ok": True,
            "job_id": job_id,
            "status": job.get("status"),
            "stage": job.get("stage"),
            "source": job.get("source"),
            "target": job.get("target"),
            "input_size": job.get("size"),
        }

        if job.get("status") == "done":
            payload.update({
                "filename": job.get("download_name"),
                "output_size": job.get("output_size"),
                "download_url": f"/api/jobs/{job_id}/download",
            })

        if job.get("status") == "error":
            payload["error"] = job.get(
                "error",
                "Conversion failed.",
            )

        return payload


@app.get("/api/jobs/{job_id}/download")
def conversion_job_download(
    job_id: str,
) -> FileResponse:
    with JOBS_LOCK:
        job = JOBS.get(job_id)

        if not job:
            raise HTTPException(
                404,
                "Conversion job was not found or has expired.",
            )

        status = job.get("status")

        if status == "error":
            raise HTTPException(
                422,
                job.get("error", "Conversion failed."),
            )

        if status != "done":
            raise HTTPException(
                409,
                "Conversion is still running.",
            )

        output_path = Path(
            str(job.get("output_path", ""))
        )
        download_name = str(
            job.get("download_name")
            or
            output_path.name
        )

    if (
        not output_path.exists()
        or
        output_path.stat().st_size == 0
    ):
        _remove_job(job_id)

        raise HTTPException(
            410,
            "The converted file is no longer available.",
        )

    cleanup = BackgroundTask(
        _remove_job,
        job_id,
    )

    return FileResponse(
        path=output_path,
        media_type=mime_for(output_path),
        filename=download_name,
        headers={
            "Cache-Control": "no-store",
        },
        background=cleanup,
    )

@app.get("/")
def service_root() -> dict[str, Any]:
    return {
        "ok": True,
        "service": "MazeDocs Converter API",
        "health": "/api",
        "max_upload_mb": MAX_UPLOAD_BYTES // (1024 * 1024),
    }
